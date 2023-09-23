import streamlit as st
import docx
import PyPDF2
from pptx import Presentation
import os
import openai
from docx import Document
from collections import defaultdict

# OpenAI API Call
def query_openai(api_key, messages):
    openai.api_key = api_key  # Utilizing the passed api_key
    response = openai.ChatCompletion.create(model="gpt-3.5-turbo", messages=messages)
    return response.choices[0].message["content"]

# Text Extraction Functions
def extract_text_from_txt(file):
    try:
        return file.getvalue().decode('utf-8')
    except Exception as e:
        st.error(f"Error processing text file. Error: {e}")
        return None

def extract_text_from_docx(file):
    try:
        doc = docx.Document(file)
        return ' '.join([para.text for para in doc.paragraphs])
    except Exception as e:
        st.error(f"Error processing docx file. Error: {e}")
        return None

def extract_text_from_pdf(file):
    try:
        pdf_reader = PyPDF2.PdfFileReader(file)
        text = ""
        for page_num in range(pdf_reader.numPages):
            text += pdf_reader.getPage(page_num).extractText()
        return text
    except Exception as e:
        st.error(f"Error processing pdf file. Error: {e}")
        return None

def extract_text_from_ppt(file):
    try:
        prs = Presentation(file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text
        return text
    except Exception as e:
        st.error(f"Error processing ppt file. Error: {e}")
        return None

def extract_text(uploaded_file):
    try:
        if uploaded_file.type == "text/plain":
            return extract_text_from_txt(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return extract_text_from_docx(uploaded_file)
        elif uploaded_file.type == "application/pdf":
            return extract_text_from_pdf(uploaded_file)
        elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            return extract_text_from_ppt(uploaded_file)
        else:
            return None
    except Exception as e:
        st.error(f"Error processing {uploaded_file.name}. Error: {e}")
        return None

# Constants
MAX_TOKENS = 200

# Streamlit Interface
st.title("Transcript Analysis Tool")

api_key = st.text_input("API Key", type="password")

uploaded_files = st.file_uploader(
    "Choose transcript files (.txt, .docx, .pdf, .ppt)",
    type=["txt", "docx", "pdf", "ppt"],
    accept_multiple_files=True,
)

guiding_questions = st.text_area("Enter the guiding questions or keywords (separated by commas)")

# Adding Streamlit caching to store the results of API calls
def extract_insights(text):
    try:
        # Trim the text if it's too long
        if len(text) > 4000:
            trimmed_text = text[:4000]  # This is a naive trim
            st.warning("The text has been truncated for analysis.")
        else:
            trimmed_text = text

        messages = [
            {"role": "system", "content": "You are a helpful assistant."},
            {
                "role": "user",
                "content": f"Summarize the following text and identify customer segments, pain points, and opportunities: {trimmed_text}"
            }
        ]
        insights = query_openai(api_key, messages)  # updated line
        return insights
    except Exception as e:
        st.error(f"OpenAI API error: {e}")



@st.cache(allow_output_mutation=True)
def generate_summary(insight):
    try:
        messages = [
            {"role": "system", "content": "You are a helpful assistant."},
            {"role": "user", "content": f"Expand on the following insight: {insight}"}
        ]
        summary = query_openai(api_key, messages)  # updated line
        return summary
    except Exception as e:
        st.error(f"OpenAI API error: {e}")

file_contents = defaultdict(str)  # updated to default dict for error handling

# Submit button action
if st.button("Submit", key='submit') and guiding_questions and uploaded_files:
    with st.spinner('Processing...'):  # Loading wheel during processing
        accepted_files = uploaded_files
        file_contents = {accepted_file.name: extract_text(accepted_file) for accepted_file in accepted_files}

        with st.expander("Consolidated Insights & Summaries"):
            for file_name, text_content in file_contents.items():
                insights = extract_insights(text_content)
                st.write(f"Insights from {file_name}:")
                st.write(insights)
                summary = generate_summary(insights)
                st.write(f"Expanded Summary for {file_name}:")
                st.write(summary)
        st.success('Processing complete!')  # Notify user of completion
        st.session_state['processing_complete'] = True  # Set session state

# Exporting Findings
def export_findings(findings_dict):
    doc = Document()
    doc.add_heading('Consolidated Findings', 0)

    for file_name, insights in findings_dict.items():
        doc.add_heading(file_name, level=1)
        doc.add_paragraph(insights)

    doc_path = 'findings.docx'
    doc.save(doc_path)
    return doc_path

if 'processing_complete' in st.session_state and st.button("Export Findings"):
    findings_dict = {file_name: extract_insights(text_content) for file_name, text_content in file_contents.items()}
    doc_path = export_findings(findings_dict)
    with open(doc_path, "rb") as file:
        st.download_button(
            label="Download Findings",
            data=file.read(),
            file_name='findings.docx',
            mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
